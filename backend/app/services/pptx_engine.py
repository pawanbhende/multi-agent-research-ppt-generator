"""
PPTX Generation Engine — Phase 1 placeholder.

NOTE: This is intentionally minimal for Phase 1, just enough to produce a
real, openable .pptx file so the full LangGraph pipeline is runnable
end-to-end today. Phase 2 will replace the slide-building internals with:
  - Custom geometric/corporate layouts per `SlideContent.layout`
  - A defined color palette + typography system
  - Automated text-wrapping / overflow protection
  - Metric callouts, comparison tables, image placeholders, etc.

The function signature below is the stable contract the Designer agent
depends on — Phase 2 work should NOT need to change this signature.
"""

import re
import uuid
from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt

from app.agents.state import SlideContent
from app.core.config import get_settings
from app.core.logging_config import logger

settings = get_settings()


def _slugify(text: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "-", text.strip().lower()).strip("-")
    return slug[:50] or "presentation"


async def generate_presentation(topic: str, outline: List[SlideContent]) -> str:
    """
    Renders a basic but real .pptx from the structured outline.
    Returns the absolute file path of the generated deck.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    for slide_data in outline:
        layout = title_layout if slide_data.layout == "title" else content_layout
        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = slide_data.title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)

        if len(slide.placeholders) > 1 and slide_data.bullets:
            body = slide.placeholders[1].text_frame
            body.clear()
            for i, bullet in enumerate(slide_data.bullets):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)

        if slide_data.speaker_notes:
            slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes

    output_dir = settings.output_path
    filename = f"{_slugify(topic)}-{uuid.uuid4().hex[:8]}.pptx"
    output_path = output_dir / filename
    prs.save(str(output_path))

    logger.info(f"[PPTX Engine] Saved {len(outline)} slides to {output_path}")
    return str(output_path.resolve())
